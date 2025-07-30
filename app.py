import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import uuid

# --- Page Configuration ---
st.set_page_config(layout="wide", page_title="Storage Bay Designer")

# --- Helper Functions ---
def hex_to_rgb(hex_color):
    """Converts a hex color string to an RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def draw_dimension_line(ax, x1, y1, x2, y2, text, is_vertical=False, offset=10, color='black', fontsize=9):
    """Draws a dimension line with arrows and text on the matplotlib axis."""
    ax.plot([x1, x2], [y1, y2], color=color, lw=1)
    if is_vertical:
        ax.plot(x1, y1, marker='v', color=color, markersize=5)
        ax.plot(x2, y2, marker='^', color=color, markersize=5)
        ax.text(x1 + offset, (y1 + y2) / 2, text, va='center', ha='left', fontsize=fontsize, rotation=90, color=color)
    else:
        ax.plot(x1, y1, marker='<', color=color, markersize=5)
        ax.plot(x2, y2, marker='>', color=color, markersize=5)
        ax.text((x1 + x2) / 2, y1 + offset, text, va='bottom', ha='center', fontsize=fontsize, color=color)

def validate_group_params(params):
    """Validates bay group parameters and returns errors if any."""
    errors = []
    if params['num_bays'] < 1:
        errors.append("Number of bays must be at least 1.")
    if params['bay_width'] <= 0:
        errors.append("Bay width must be positive.")
    if params['total_height'] <= 0:
        errors.append("Total height must be positive.")
    if params['ground_clearance'] < 0:
        errors.append("Ground clearance cannot be negative.")
    if params['shelf_thickness'] <= 0:
        errors.append("Shelf thickness must be positive.")
    if params['side_panel_thickness'] <= 0:
        errors.append("Side panel thickness must be positive.")
    if params.get('bin_split_thickness', 0) <= 0:
        errors.append("Bin split thickness must be positive.")
    if params['num_cols'] < 1:
        errors.append("Number of columns must be at least 1.")
    if params['num_rows'] < 1:
        errors.append("Number of rows must be at least 1.")
    total_net_bin_h = sum(params['bin_heights'])
    num_shelves = params['num_rows'] + (1 if params['has_top_cap'] else 0)
    required_height = total_net_bin_h + num_shelves * params['shelf_thickness'] + params['ground_clearance']
    if abs(required_height - params['total_height']) > 0.1:
        errors.append(f"Calculated height ({required_height:.1f} mm) does not match target height ({params['total_height']:.1f} mm).")
    return errors

@st.cache_data
def draw_bay_group(params):
    """Main function to draw a group of bays using Matplotlib for the LIVE PREVIEW."""
    # Unpack parameters
    num_bays = params['num_bays']
    bay_width = params['bay_width']
    total_height = params['total_height']
    ground_clearance = params['ground_clearance']
    shelf_thickness = params['shelf_thickness']
    side_panel_thickness = params['side_panel_thickness']
    bin_split_thickness = params['bin_split_thickness']
    num_cols = params['num_cols']
    num_rows = params['num_rows']
    has_top_cap = params['has_top_cap']
    color = params['color']
    bin_heights = params['bin_heights']
    zoom_factor = params.get('zoom', 1.0)

    # Normalize visual thickness to prevent bulky appearance
    visual_shelf_thickness = min(shelf_thickness, 18.0)  # Cap at 18 mm for visual rendering
    visual_bin_split_thickness = min(bin_split_thickness, 18.0)  # Cap at 18 mm for visual rendering
    visual_side_panel_thickness = max(side_panel_thickness, 10.0)  # Minimum 10 mm for visual rendering

    # --- Calculations ---
    core_width = num_bays * bay_width
    total_group_width = core_width + (2 * side_panel_thickness)  # Use actual thickness for calculations
    dim_offset_x = 0.05 * core_width
    dim_offset_y = 0.05 * total_height
    
    fig, ax = plt.subplots(figsize=(12, 12))

    # --- Draw Side Panels ---
    def draw_side_panels():
        ax.add_patch(patches.Rectangle((-visual_side_panel_thickness, 0), visual_side_panel_thickness, total_height, facecolor='none', edgecolor=color, lw=1))
        ax.add_patch(patches.Rectangle((core_width, 0), visual_side_panel_thickness, total_height, facecolor='none', edgecolor=color, lw=1))

    # --- Draw Bays ---
    def draw_bays():
        current_x = 0
        for bay_idx in range(num_bays):
            net_width_per_bay = bay_width - 2 * side_panel_thickness  # Account for both side panels
            total_internal_dividers = (num_cols - 1) * bin_split_thickness
            bin_width = (net_width_per_bay - total_internal_dividers) / num_cols if num_cols > 0 else 0

            bin_start_x = current_x
            if num_cols > 1:
                for i in range(1, num_cols):
                    split_x = bin_start_x + (i * bin_width) + ((i-1) * bin_split_thickness)
                    ax.add_patch(patches.Rectangle((split_x, ground_clearance), visual_bin_split_thickness, structure_height, facecolor='none', edgecolor=color, lw=1))
            
            if bay_idx < num_bays - 1:
                divider_x = current_x + bay_width
                ax.plot([divider_x, divider_x], [ground_clearance, structure_height], color='#aaaaaa', lw=1, linestyle='--')

            current_x += bay_width

    # --- Draw Shelves and Dimensions ---
    def draw_shelves_and_dimensions():
        current_y = ground_clearance
        pitch_offset_x = dim_offset_x * 2.5

        for i in range(num_rows):
            shelf_bottom_y = current_y
            ax.add_patch(patches.Rectangle((-visual_side_panel_thickness, shelf_bottom_y), total_group_width, visual_shelf_thickness, facecolor='none', edgecolor=color, lw=1))
            shelf_top_y = shelf_bottom_y + shelf_thickness  # Use actual thickness for positioning
            
            if i < len(bin_heights):
                net_bin_h = bin_heights[i]
                pitch_h = net_bin_h + shelf_thickness
                level_name = chr(65 + i)
                
                bin_bottom_y = shelf_top_y
                bin_top_y = bin_bottom_y + net_bin_h
                draw_dimension_line(ax, core_width + visual_side_panel_thickness + dim_offset_x, bin_bottom_y, core_width + visual_side_panel_thickness + dim_offset_x, bin_top_y, f"{net_bin_h:.1f}", is_vertical=True, offset=5, color='#3b82f6')
                
                pitch_top_y = shelf_bottom_y + pitch_h
                draw_dimension_line(ax, core_width + visual_side_panel_thickness + pitch_offset_x, shelf_bottom_y, core_width + visual_side_panel_thickness + pitch_offset_x, pitch_top_y, f"{pitch_h:.1f}", is_vertical=True, offset=5, color='black')

                ax.text(-visual_side_panel_thickness - dim_offset_x, (bin_bottom_y + bin_top_y) / 2, level_name, va='center', ha='center', fontsize=12, fontweight='bold')
                
                current_y = bin_top_y

        if has_top_cap:
            ax.add_patch(patches.Rectangle((-visual_side_panel_thickness, total_height - visual_shelf_thickness), total_group_width, visual_shelf_thickness, facecolor='none', edgecolor=color, lw=1))

    # --- Draw Main Dimensions ---
    def draw_main_dimensions():
        draw_dimension_line(ax, -visual_side_panel_thickness, -dim_offset_y * 2, core_width + visual_side_panel_thickness, -dim_offset_y * 2, f"Total Group Width: {total_group_width:.0f} mm", offset=10)
        draw_dimension_line(ax, -visual_side_panel_thickness - (dim_offset_x * 4), 0, -visual_side_panel_thickness - (dim_offset_x * 4), total_height, f"Total Height: {total_height:.0f} mm", is_vertical=True, offset=10)

        if num_cols > 0:
            dim_y_pos = total_height + dim_offset_y
            loop_current_x = 0
            for bay_idx in range(num_bays):
                net_width_per_bay = bay_width - 2 * side_panel_thickness  # Account for both side panels
                total_internal_dividers = (num_cols - 1) * bin_split_thickness
                bin_width = (net_width_per_bay - total_internal_dividers) / num_cols if num_cols > 0 else 0
                
                bin_start_x = loop_current_x
                for i in range(num_cols):
                    dim_start_x = bin_start_x + (i * (bin_width + bin_split_thickness))
                    dim_end_x = dim_start_x + bin_width
                    draw_dimension_line(ax, dim_start_x, dim_y_pos, dim_end_x, dim_y_pos, f"{bin_width:.1f}", offset=10, color='#3b82f6')
                
                loop_current_x += bay_width

    # --- Execute Drawing ---
    structure_height = total_height - ground_clearance
    draw_side_panels()
    draw_bays()
    draw_shelves_and_dimensions()
    draw_main_dimensions()

    # --- Final Touches ---
    ax.set_aspect('equal', adjustable='box')
    padding_x = core_width * 0.4 + visual_side_panel_thickness
    ax.set_xlim((-padding_x) * zoom_factor, (core_width + padding_x) * zoom_factor)
    ax.set_ylim(-dim_offset_y * 3 * zoom_factor, total_height + dim_offset_y * 2 * zoom_factor)
    ax.axis('off')
    
    return fig

def create_summary_slide(prs, bay_groups):
    """Creates a summary slide with a bill of materials."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_shape.text = "Bill of Materials"
    
    table_left, table_top, table_width, table_height = Inches(1), Inches(1), Inches(8), Inches(4)
    table = slide.shapes.add_table(2 + len(bay_groups), 4, table_left, table_top, table_width, table_height).table
    
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2)
    
    table.cell(0, 0).text = "Group Name"
    table.cell(0, 1).text = "Side Panels"
    table.cell(0, 2).text = "Shelves"
    table.cell(0, 3).text = "Bin Dividers"
    
    total_side_panels = 0
    total_shelves = 0
    total_dividers = 0
    
    for i, group in enumerate(bay_groups, 1):
        side_panels = 2
        shelves = group['num_rows'] + (1 if group['has_top_cap'] else 0)
        dividers = (group['num_cols'] - 1) * group['num_bays']
        
        table.cell(i, 0).text = group['name']
        table.cell(i, 1).text = str(side_panels)
        table.cell(i, 2).text = str(shelves)
        table.cell(i, 3).text = str(dividers)
        
        total_side_panels += side_panels
        total_shelves += shelves
        total_dividers += dividers
    
    table.cell(len(bay_groups) + 1, 0).text = "Total"
    table.cell(len(bay_groups) + 1, 1).text = str(total_side_panels)
    table.cell(len(bay_groups) + 1, 2).text = str(total_shelves)
    table.cell(len(bay_groups) + 1, 3).text = str(total_dividers)

def create_editable_powerpoint(bay_groups):
    """Creates a PowerPoint presentation from bay group data using native shapes."""
    prs = Presentation()
    
    create_summary_slide(prs, bay_groups)
    
    for group_data in bay_groups:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_shape.text = f"Design for: {group_data['name']}"

        # --- Unpack Parameters ---
        num_bays, bay_width, total_height, ground_clearance, shelf_thickness, side_panel_thickness, bin_split_thickness, num_cols, num_rows, has_top_cap, color_hex, bin_heights = (
            group_data['num_bays'], group_data['bay_width'], group_data['total_height'],
            group_data['ground_clearance'], group_data['shelf_thickness'], group_data['side_panel_thickness'],
            group_data['bin_split_thickness'], group_data['num_cols'], group_data['num_rows'], group_data['has_top_cap'],
            group_data['color'], group_data['bin_heights']
        )

        # Normalize visual thickness for PowerPoint
        visual_shelf_thickness = min(shelf_thickness, 18.0)  # Cap at 18 mm for visual rendering
        visual_bin_split_thickness = min(bin_split_thickness, 18.0)  # Cap at 18 mm for visual rendering
        visual_side_panel_thickness = max(side_panel_thickness, 10.0)  # Minimum 10 mm for visual rendering

        # --- Define Drawing Area and Scale on Slide ---
        canvas_left, canvas_top, canvas_width, canvas_height = Inches(1.5), Inches(1), Inches(7), Inches(5.5)
        total_group_width = (num_bays * bay_width) + (2 * side_panel_thickness)  # Use actual thickness for calculations
        scale = min(max(min(canvas_width / total_group_width, canvas_height / total_height), 0.1), 1.0)  # Increased minimum scale to 0.1

        def pt_to_emu(points):
            return int(points * 12700)

        def add_shape(left_mm, top_mm, width_mm, height_mm, color_hex):
            if left_mm < 0 or top_mm < 0 or width_mm <= 0 or height_mm <= 0 or scale <= 0:
                st.error(f"Invalid shape parameters: left={left_mm}, top={top_mm}, width={width_mm}, height={height_mm}, scale={scale}")
                return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, canvas_left, canvas_top, Inches(0.1), Inches(0.1))  # Fallback shape
            left = canvas_left + left_mm * scale
            top = canvas_top + (top_mm) * scale  # Bottom-up alignment
            width = max(width_mm * scale, Inches(0.05))
            height = max(height_mm * scale, Inches(0.05))
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.fill.solid()  # Set to no fill
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill to simulate no fill
            shape.line.color.rgb = RGBColor(*hex_to_rgb(color_hex))  # Use selected color for border
            shape.line.width = Pt(0.5)  # Thin border to mimic sticks
            return shape
        
        def add_dimension(start_x, start_y, end_x, end_y, text, is_vertical=False):
            if start_x >= end_x or start_y >= end_y or scale <= 0:
                st.error(f"Invalid dimension parameters: start_x={start_x}, start_y={start_y}, end_x={end_x}, end_y={end_y}, scale={scale}")
                return
            line = slide.shapes.add_connector(1, start_x, start_y, end_x, end_y)
            line.line.fill.solid()
            line.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
            line.line.begin_arrow_type = 2
            line.line.end_arrow_type = 2

            if is_vertical:
                text_left = start_x + pt_to_emu(5)
                text_top = start_y + (end_y - start_y) / 2 - pt_to_emu(20)
                textbox = slide.shapes.add_textbox(text_left, text_top, Inches(0.5), Inches(0.5))
                textbox.rotation = 270.0
            else:
                text_left = start_x + (end_x - start_x) / 2 - pt_to_emu(20)
                text_top = start_y - pt_to_emu(12)
                textbox = slide.shapes.add_textbox(text_left, text_top, Inches(0.5), Inches(0.5))
            
            p = textbox.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(8)
            p.alignment = PP_ALIGN.CENTER
            textbox.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # --- Draw Structure ---
        structure_height = total_height - ground_clearance
        shape = add_shape(0, 0, visual_side_panel_thickness, total_height, color_hex)
        current_x_mm = visual_side_panel_thickness
        if current_x_mm > total_group_width:
            st.error(f"current_x_mm {current_x_mm} exceeds total_group_width {total_group_width}")

        for bay_idx in range(num_bays):
            net_width_per_bay = bay_width - 2 * side_panel_thickness  # Account for both side panels
            total_internal_dividers = (num_cols - 1) * bin_split_thickness
            bin_width = (net_width_per_bay - total_internal_dividers) / num_cols if num_cols > 0 else 0

            bin_start_x_mm = current_x_mm
            if num_cols > 1:
                for i in range(1, num_cols):
                    split_x_mm = bin_start_x_mm + (i * bin_width) + ((i-1) * bin_split_thickness)
                    if split_x_mm > total_group_width:
                        st.error(f"split_x_mm {split_x_mm} exceeds total_group_width {total_group_width}")
                    add_shape(split_x_mm, ground_clearance, visual_bin_split_thickness, structure_height, color_hex)
            
            # Add inner bin width dimensions
            for i in range(num_cols):
                dim_start_x = canvas_left + (bin_start_x_mm + i * (bin_width + bin_split_thickness)) * scale
                dim_end_x = dim_start_x + (bin_width * scale)
                dim_y = canvas_top - pt_to_emu(20)
                add_dimension(dim_start_x, dim_y, dim_end_x, dim_y, f"{bin_width:.1f} mm")

            current_x_mm += bay_width
            if current_x_mm > total_group_width:
                st.error(f"current_x_mm {current_x_mm} exceeds total_group_width {total_group_width} after bay {bay_idx}")

        add_shape(current_x_mm, 0, visual_side_panel_thickness, total_height, color_hex)

        current_y_mm = ground_clearance
        if current_y_mm > total_height:
            st.error(f"current_y_mm {current_y_mm} exceeds total_height {total_height}")

        for i in range(num_rows):
            shelf_bottom_y = current_y_mm
            add_shape(0, shelf_bottom_y, total_group_width, visual_shelf_thickness, color_hex)
            shelf_top_y = shelf_bottom_y + shelf_thickness  # Use actual thickness for positioning
            
            if i < len(bin_heights):
                net_bin_h = bin_heights[i]
                pitch_h = net_bin_h + shelf_thickness

                dim_start_y = canvas_top + (shelf_top_y) * scale
                dim_end_y = canvas_top + (shelf_top_y + net_bin_h) * scale
                dim_x = canvas_left + (total_group_width + 50) * scale
                add_dimension(dim_x, dim_start_y, dim_x, dim_end_y, f"{net_bin_h:.1f} mm", is_vertical=True)
                
                pitch_dim_start_y = canvas_top + (shelf_bottom_y) * scale
                pitch_dim_end_y = canvas_top + (shelf_bottom_y + pitch_h) * scale
                pitch_dim_x = canvas_left + (total_group_width + 150) * scale
                add_dimension(pitch_dim_x, pitch_dim_start_y, pitch_dim_x, pitch_dim_end_y, f"{pitch_h:.1f} mm", is_vertical=True)

                current_y_mm += shelf_thickness + net_bin_h
                if current_y_mm > total_height:
                    st.error(f"current_y_mm {current_y_mm} exceeds total_height {total_height} at row {i}")

        if has_top_cap:
            add_shape(0, total_height - visual_shelf_thickness, total_group_width, visual_shelf_thickness, color_hex)

        # Add outer dimensions
        total_w_y = canvas_top + canvas_height + pt_to_emu(20)
        add_dimension(canvas_left, total_w_y, canvas_left + total_group_width * scale, total_w_y, f"Total Width: {total_group_width:.0f} mm")
        
        total_h_x = canvas_left - pt_to_emu(40)
        add_dimension(total_h_x, canvas_top, total_h_x, canvas_top + total_height * scale, f"Total Height: {total_height:.0f} mm", is_vertical=True)

    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)  # Ensure cursor is at start
    if ppt_buf.getbuffer().nbytes == 0:
        st.error("Generated PPTX file is empty. Please check configuration and try again.")
    return ppt_buf

# --- Initialize Session State ---
if 'bay_groups' not in st.session_state:
    st.session_state.bay_groups = [{
        "id": str(uuid.uuid4()),
        "name": "Group A",
        "num_bays": 2,
        "bay_width": 1050.0,
        "total_height": 2000.0,
        "ground_clearance": 50.0,
        "shelf_thickness": 18.0,
        "side_panel_thickness": 18.0,
        "bin_split_thickness": 18.0,
        "num_cols": 4,
        "num_rows": 5,
        "has_top_cap": True,
        "color": "#4A90E2",
        "bin_heights": [350.0] * 5,
        "lock_heights": [False] * 5
    }]

# Migrate existing groups to include bin_split_thickness
for group in st.session_state.bay_groups:
    if 'bin_split_thickness' not in group:
        group['bin_split_thickness'] = 18.0  # Default value for existing groups

# --- Sidebar Controls ---
st.sidebar.header("Manage Bay Groups")

with st.sidebar.expander("Add New Group", expanded=True):
    with st.form("new_group_form"):
        new_group_name = st.text_input("New Group Name", "New Group", help="Enter a unique name for the new bay group.")
        add_group_submitted = st.form_submit_button("Add Group")
        if add_group_submitted:
            if any(g['name'] == new_group_name for g in st.session_state.bay_groups):
                st.error("Group name must be unique.")
            else:
                new_group = st.session_state.bay_groups[0].copy()
                new_group['id'] = str(uuid.uuid4())
                new_group['name'] = new_group_name
                new_group['bin_split_thickness'] = 18.0  # Ensure new groups have bin_split_thickness
                st.session_state.bay_groups.append(new_group)
                st.rerun()

if len(st.session_state.bay_groups) > 1:
    if st.sidebar.button("Remove Last Group", help="Removes the most recently added group."):
        st.session_state.bay_groups.pop()
        st.rerun()

st.sidebar.markdown("---")

group_names = [g['name'] for g in st.session_state.bay_groups]
selected_group_name = st.sidebar.selectbox("Select Group to Edit", group_names, help="Choose a group to modify its settings.")
active_group_idx = group_names.index(selected_group_name)
group_data = st.session_state.bay_groups[active_group_idx]

# --- Dynamic Height Calculation Callbacks ---
def distribute_total_height():
    active_group = st.session_state.bay_groups[active_group_idx]
    num_shelves_for_calc = active_group['num_rows'] + (1 if active_group['has_top_cap'] else 0)
    total_shelf_thickness = num_shelves_for_calc * active_group['shelf_thickness']
    available_space = active_group['total_height'] - active_group['ground_clearance'] - total_shelf_thickness
    
    unlocked_indices = [i for i, locked in enumerate(active_group['lock_heights']) if not locked]
    num_unlocked = len(unlocked_indices)
    
    if available_space > 0 and num_unlocked > 0:
        uniform_net_h = available_space / num_unlocked
        for i in unlocked_indices:
            active_group['bin_heights'][i] = uniform_net_h

def update_total_height():
    """Update total_height based on bin heights, shelf thickness, and ground clearance."""
    active_group = st.session_state.bay_groups[active_group_idx]
    num_shelves_for_calc = active_group['num_rows'] + (1 if active_group['has_top_cap'] else 0)
    total_shelf_thickness = num_shelves_for_calc * active_group['shelf_thickness']
    total_net_bin_h = sum(active_group['bin_heights'])
    active_group['total_height'] = total_net_bin_h + total_shelf_thickness + active_group['ground_clearance']

# --- Configuration Inputs ---
with st.sidebar.expander("Structure", expanded=True):
    group_data['num_bays'] = st.number_input("Number of Bays in Group", min_value=1, value=int(group_data['num_bays']), key=f"num_bays_{group_data['id']}", help="Number of bays in the group.")
    group_data['bay_width'] = st.number_input("Width per Bay (mm)", min_value=1.0, value=float(group_data['bay_width']), key=f"bay_width_{group_data['id']}", help="Width of each bay in millimeters.")
    group_data['total_height'] = st.number_input("Target Total Height (mm)", min_value=1.0, value=float(group_data['total_height']), key=f"total_height_{group_data['id']}", on_change=distribute_total_height, help="Set to automatically distribute height among unlocked bins.")
    group_data['ground_clearance'] = st.number_input("Ground Clearance (mm)", min_value=0.0, value=float(group_data['ground_clearance']), key=f"ground_clearance_{group_data['id']}", on_change=update_total_height, help="Height from ground to first shelf.")
    group_data['has_top_cap'] = st.checkbox("Add Top Cap", value=group_data['has_top_cap'], key=f"has_top_cap_{group_data['id']}", on_change=update_total_height, help="Include a top cap shelf.")

with st.sidebar.expander("Layout", expanded=True):
    prev_num_rows = group_data['num_rows']
    group_data['num_rows'] = st.number_input("Shelves (Rows)", min_value=1, value=int(group_data['num_rows']), key=f"num_rows_{group_data['id']}", on_change=update_total_height, help="Number of horizontal shelves.")
    group_data['num_cols'] = st.number_input("Bin-Split (Columns)", min_value=1, value=int(group_data['num_cols']), key=f"num_cols_{group_data['id']}", help="Number of vertical bin splits per bay.")

    # Adjust bin_heights and lock_heights if num_rows changes
    if prev_num_rows != group_data['num_rows']:
        if group_data['num_rows'] > len(group_data['bin_heights']):
            # Extend with default values
            default_height = group_data['bin_heights'][0] if group_data['bin_heights'] else 350.0
            group_data['bin_heights'].extend([default_height] * (group_data['num_rows'] - len(group_data['bin_heights'])))
            group_data['lock_heights'].extend([False] * (group_data['num_rows'] - len(group_data['lock_heights'])))
        else:
            # Trim excess
            group_data['bin_heights'] = group_data['bin_heights'][:group_data['num_rows']]
            group_data['lock_heights'] = group_data['lock_heights'][:group_data['num_rows']]
        update_total_height()

with st.sidebar.expander("Individual Net Bin Heights", expanded=True):
    auto_distribute = st.checkbox("Auto-distribute Heights", value=True, key=f"auto_distribute_{group_data['id']}", help="Automatically distribute heights among unlocked bins.")
    
    current_bin_heights = []
    current_lock_heights = []
    for j in range(group_data['num_rows']):
        level_name = chr(65 + j)
        col1, col2 = st.columns([3, 1])
        with col1:
            height = st.number_input(
                f"Level {level_name} Net Height",
                min_value=1.0,
                value=float(group_data['bin_heights'][j]),
                key=f"level_{group_data['id']}_{j}",
                disabled=auto_distribute and not group_data['lock_heights'][j],
                on_change=update_total_height
            )
        with col2:
            locked = st.checkbox("Lock", value=group_data['lock_heights'][j], key=f"lock_{group_data['id']}_{j}", help="Lock this height to prevent auto-distribution.", on_change=update_total_height)
        current_bin_heights.append(height)
        current_lock_heights.append(locked)
    
    group_data['bin_heights'] = current_bin_heights
    group_data['lock_heights'] = current_lock_heights
    if auto_distribute:
        distribute_total_height()
    else:
        update_total_height()

with st.sidebar.expander("Materials & Appearance", expanded=True):
    group_data['shelf_thickness'] = st.number_input(
        "Shelf Thickness (mm)", 
        min_value=1.0, 
        value=float(group_data['shelf_thickness']), 
        key=f"shelf_thick_{group_data['id']}", 
        on_change=update_total_height, 
        help="Thickness of horizontal shelves. Large values are used in calculations but rendered as 18 mm max for visual clarity."
    )
    group_data['bin_split_thickness'] = st.number_input(
        "Bin Split Thickness (mm)", 
        min_value=1.0, 
        value=float(group_data.get('bin_split_thickness', 18.0)), 
        key=f"bin_split_thick_{group_data['id']}", 
        help="Thickness of vertical bin dividers. Large values are used in calculations but rendered as 18 mm max for visual clarity."
    )
    group_data['side_panel_thickness'] = st.number_input(
        "Outer Side Panel Thickness (mm)", 
        min_value=1.0, 
        value=float(group_data['side_panel_thickness']), 
        key=f"side_panel_thick_{group_data['id']}", 
        help="Thickness of side panels. A minimum visual thickness of 10 mm is used for rendering to ensure visibility."
    )
    group_data['color'] = st.color_picker("Structure Color", value=group_data['color'], key=f"color_{group_data['id']}", help="Color of the structure.")
    group_data['zoom'] = st.slider("Zoom", 1.0, 5.0, group_data.get('zoom', 1.0), 0.1, key=f"zoom_{group_data['id']}", help="Adjust zoom level for the preview.")

# --- Validate Parameters ---
errors = validate_group_params(group_data)
if errors:
    st.sidebar.error("Configuration Errors:\n" + "\n".join(f"- {e}" for e in errors))

# --- Calculate and Display Final Height ---
total_net_bin_h = sum(group_data['bin_heights'])
num_shelves_for_calc = group_data['num_rows'] + (1 if group_data['has_top_cap'] else 0)
total_shelf_h = num_shelves_for_calc * group_data['shelf_thickness']
calculated_total_height = total_net_bin_h + total_shelf_h + group_data['ground_clearance']
st.sidebar.metric("Calculated Total Height", f"{calculated_total_height:.1f} mm")

# --- Main Area for Drawing ---
st.header(f"Generated Design for: {group_data['name']}")
if not errors:
    fig = draw_bay_group(group_data)
    st.pyplot(fig, use_container_width=True)
else:
    st.error("Please fix configuration errors to view the design.")

# --- Global Download Button ---
st.sidebar.markdown("---")
st.sidebar.header("Download All Designs")

download_button_placeholder = st.sidebar.empty()

if st.sidebar.button("Generate PPTX", help="Generate a PowerPoint file with all designs and a bill of materials."):
    has_errors = False
    for group in st.session_state.bay_groups:
        if validate_group_params(group):
            has_errors = True
            st.error(f"Cannot generate PPTX due to errors in group: {group['name']}")
    
    if not has_errors:
        ppt_buffer = create_editable_powerpoint(st.session_state.bay_groups)
        if ppt_buffer is not None:
            download_button_placeholder.download_button(
                label="Download Now",
                data=ppt_buffer.getvalue(),  # Use getvalue() to ensure full buffer content
                file_name="all_bay_designs.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                help="Download the PowerPoint file containing all bay designs."
            )
        else:
            st.error("Failed to generate PPTX file. Please check the Streamlit logs for details.")